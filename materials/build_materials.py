from pathlib import Path

from openpyxl import Workbook
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
XLSX_PATH = ROOT / "financial_model.xlsx"
PPTX_PATH = ROOT / "pitch_deck.pptx"


ASSUMPTIONS = {
    "hardware_price": 1800,
    "software_fee": 85,
    "onboarding_fee": 7500,
    "hardware_cogs": 1050,
    "cloud_cost": 12,
    "support_cost": 8,
    "eng_monthly": 18000,
    "sales_monthly": 7000,
    "tools_monthly": 2500,
    "marketing_monthly": 1500,
    "prototype_capex": 45000,
    "smart_ppe_market_2025_musd": 3300,
    "smart_ppe_market_2033_musd": 9800,
    "smart_helmet_market_2024_musd": 892,
    "smart_helmet_market_2030_musd": 2300,
    "target_orgs": 120,
    "workers_per_org": 200,
    "year2_penetration": 0.03,
    "year3_penetration": 0.07,
    "year2_new_customers": 6,
    "year3_new_customers": 10,
    "year2_opex_growth": 0.25,
    "year3_opex_growth": 0.20,
}

DEPLOYMENT = {
    "new_customers": [0, 0, 0, 0, 0, 1, 0, 1, 0, 1, 0, 1],
    "new_units": [0, 0, 0, 0, 0, 20, 0, 30, 0, 40, 0, 60],
    "active_users": [0, 0, 0, 0, 0, 20, 20, 50, 50, 90, 90, 150],
}


def style_header(ws, row=1):
    fill = PatternFill("solid", fgColor="1F4E78")
    font = Font(color="FFFFFF", bold=True)
    for cell in ws[row]:
        cell.fill = fill
        cell.font = font
        cell.alignment = Alignment(horizontal="center", vertical="center")


def autosize(ws, min_width=12, max_width=32):
    for idx, column_cells in enumerate(ws.columns, start=1):
        values = [str(cell.value) for cell in column_cells if cell.value is not None]
        if not values:
            width = min_width
        else:
            width = max(min(max(len(v) for v in values) + 2, min_width), max_width)
        ws.column_dimensions[get_column_letter(idx)].width = width


def money(value):
    if value >= 1_000_000:
        return f"C${value / 1_000_000:.2f}M"
    if value >= 1_000:
        return f"C${value / 1_000:.0f}K"
    return f"C${value:,.0f}"


def percent(value):
    return f"{value * 100:.1f}%"


def cagr(start, end, years):
    return (end / start) ** (1 / years) - 1


def calculate_summary():
    beachhead_users = ASSUMPTIONS["target_orgs"] * ASSUMPTIONS["workers_per_org"]
    ppe_cagr = cagr(
        ASSUMPTIONS["smart_ppe_market_2025_musd"],
        ASSUMPTIONS["smart_ppe_market_2033_musd"],
        8,
    )
    helmet_cagr = cagr(
        ASSUMPTIONS["smart_helmet_market_2024_musd"],
        ASSUMPTIONS["smart_helmet_market_2030_musd"],
        6,
    )

    annual_opex_y1 = (
        ASSUMPTIONS["eng_monthly"]
        + ASSUMPTIONS["sales_monthly"]
        + ASSUMPTIONS["tools_monthly"]
        + ASSUMPTIONS["marketing_monthly"]
    ) * 12

    y1_new_customers = sum(DEPLOYMENT["new_customers"])
    y1_units = sum(DEPLOYMENT["new_units"])
    y1_user_months = sum(DEPLOYMENT["active_users"])
    y1_avg_active = y1_user_months / 12
    y1_year_end_users = DEPLOYMENT["active_users"][-1]
    y1_revenue = (
        y1_units * ASSUMPTIONS["hardware_price"]
        + y1_user_months * ASSUMPTIONS["software_fee"]
        + y1_new_customers * ASSUMPTIONS["onboarding_fee"]
    )
    y1_cogs = (
        y1_units * ASSUMPTIONS["hardware_cogs"]
        + y1_user_months * ASSUMPTIONS["cloud_cost"]
        + y1_user_months * ASSUMPTIONS["support_cost"]
    )
    y1_gross_profit = y1_revenue - y1_cogs
    y1_gross_margin = y1_gross_profit / y1_revenue
    y1_arr = y1_year_end_users * ASSUMPTIONS["software_fee"] * 12
    y1_ebitda = y1_gross_profit - annual_opex_y1
    y1_net_cash = y1_ebitda - ASSUMPTIONS["prototype_capex"]

    y2_year_end_users = round(beachhead_users * ASSUMPTIONS["year2_penetration"])
    y3_year_end_users = round(beachhead_users * ASSUMPTIONS["year3_penetration"])
    y2_new_units = y2_year_end_users - y1_year_end_users
    y3_new_units = y3_year_end_users - y2_year_end_users
    y2_avg_active = (y1_year_end_users + y2_year_end_users) / 2
    y3_avg_active = (y2_year_end_users + y3_year_end_users) / 2

    annual_opex_y2 = annual_opex_y1 * (1 + ASSUMPTIONS["year2_opex_growth"])
    annual_opex_y3 = annual_opex_y2 * (1 + ASSUMPTIONS["year3_opex_growth"])

    y2_revenue = (
        y2_new_units * ASSUMPTIONS["hardware_price"]
        + y2_avg_active * 12 * ASSUMPTIONS["software_fee"]
        + ASSUMPTIONS["year2_new_customers"] * ASSUMPTIONS["onboarding_fee"]
    )
    y3_revenue = (
        y3_new_units * ASSUMPTIONS["hardware_price"]
        + y3_avg_active * 12 * ASSUMPTIONS["software_fee"]
        + ASSUMPTIONS["year3_new_customers"] * ASSUMPTIONS["onboarding_fee"]
    )
    y2_cogs = (
        y2_new_units * ASSUMPTIONS["hardware_cogs"]
        + y2_avg_active * 12 * ASSUMPTIONS["cloud_cost"]
        + y2_avg_active * 12 * ASSUMPTIONS["support_cost"]
    )
    y3_cogs = (
        y3_new_units * ASSUMPTIONS["hardware_cogs"]
        + y3_avg_active * 12 * ASSUMPTIONS["cloud_cost"]
        + y3_avg_active * 12 * ASSUMPTIONS["support_cost"]
    )
    y2_gross_profit = y2_revenue - y2_cogs
    y3_gross_profit = y3_revenue - y3_cogs
    y2_gross_margin = y2_gross_profit / y2_revenue
    y3_gross_margin = y3_gross_profit / y3_revenue
    y2_arr = y2_year_end_users * ASSUMPTIONS["software_fee"] * 12
    y3_arr = y3_year_end_users * ASSUMPTIONS["software_fee"] * 12

    return {
        "beachhead_users": beachhead_users,
        "ppe_cagr": ppe_cagr,
        "helmet_cagr": helmet_cagr,
        "annual_opex_y1": annual_opex_y1,
        "y1": {
            "new_customers": y1_new_customers,
            "units": y1_units,
            "user_months": y1_user_months,
            "avg_active": y1_avg_active,
            "year_end_users": y1_year_end_users,
            "revenue": y1_revenue,
            "cogs": y1_cogs,
            "gross_profit": y1_gross_profit,
            "gross_margin": y1_gross_margin,
            "arr": y1_arr,
            "ebitda": y1_ebitda,
            "net_cash": y1_net_cash,
        },
        "y2": {
            "new_customers": ASSUMPTIONS["year2_new_customers"],
            "units": y2_new_units,
            "avg_active": y2_avg_active,
            "year_end_users": y2_year_end_users,
            "revenue": y2_revenue,
            "cogs": y2_cogs,
            "gross_profit": y2_gross_profit,
            "gross_margin": y2_gross_margin,
            "arr": y2_arr,
            "opex": annual_opex_y2,
            "ebitda": y2_gross_profit - annual_opex_y2,
        },
        "y3": {
            "new_customers": ASSUMPTIONS["year3_new_customers"],
            "units": y3_new_units,
            "avg_active": y3_avg_active,
            "year_end_users": y3_year_end_users,
            "revenue": y3_revenue,
            "cogs": y3_cogs,
            "gross_profit": y3_gross_profit,
            "gross_margin": y3_gross_margin,
            "arr": y3_arr,
            "opex": annual_opex_y3,
            "ebitda": y3_gross_profit - annual_opex_y3,
        },
    }


def build_financial_model():
    summary = calculate_summary()
    wb = Workbook()
    dashboard = wb.active
    dashboard.title = "Dashboard"

    ws = wb.create_sheet("Assumptions")
    ws.title = "Assumptions"

    ws.append(["Category", "Variable", "Value", "Notes"])
    assumption_rows = [
        ("hardware_price", "Pricing", "Hardware price per unit", ASSUMPTIONS["hardware_price"], "Pilot-stage sale price in CAD"),
        ("software_fee", "Pricing", "Monthly software fee per active user", ASSUMPTIONS["software_fee"], "Per-user SaaS fee"),
        ("onboarding_fee", "Pricing", "Onboarding fee per new customer", ASSUMPTIONS["onboarding_fee"], "Setup and training"),
        ("hardware_cogs", "COGS", "Hardware COGS per unit", ASSUMPTIONS["hardware_cogs"], "Electronics, assembly, shipping"),
        ("cloud_cost", "COGS", "Cloud cost per active user per month", ASSUMPTIONS["cloud_cost"], "Inference, storage, compute"),
        ("support_cost", "COGS", "Support cost per active user per month", ASSUMPTIONS["support_cost"], "Customer support and success"),
        ("eng_monthly", "OPEX", "Engineering payroll per month", ASSUMPTIONS["eng_monthly"], "ML plus embedded/product contractor mix"),
        ("sales_monthly", "OPEX", "Sales and business development per month", ASSUMPTIONS["sales_monthly"], "Founder sales plus travel"),
        ("tools_monthly", "OPEX", "Software/tools/admin per month", ASSUMPTIONS["tools_monthly"], "SaaS, insurance, legal, admin"),
        ("marketing_monthly", "OPEX", "Marketing/events per month", ASSUMPTIONS["marketing_monthly"], "Industry events and outreach"),
        ("prototype_capex", "CAPEX", "Prototype and test equipment", ASSUMPTIONS["prototype_capex"], "Spread across the first 4 months"),
        ("smart_ppe_market_2025_musd", "Market", "Global smart PPE market size, 2025 (USD M)", ASSUMPTIONS["smart_ppe_market_2025_musd"], "From team status report market citation"),
        ("smart_ppe_market_2033_musd", "Market", "Global smart PPE market size, 2033 (USD M)", ASSUMPTIONS["smart_ppe_market_2033_musd"], "From team status report market citation"),
        ("smart_helmet_market_2024_musd", "Market", "Global smart helmet market size, 2024 (USD M)", ASSUMPTIONS["smart_helmet_market_2024_musd"], "From team status report market citation"),
        ("smart_helmet_market_2030_musd", "Market", "Global smart helmet market size, 2030 (USD M)", ASSUMPTIONS["smart_helmet_market_2030_musd"], "From team status report market citation"),
        ("target_orgs", "Market", "Target organizations in beachhead", ASSUMPTIONS["target_orgs"], "Planning assumption to be validated"),
        ("workers_per_org", "Market", "Average qualifying workers per organization", ASSUMPTIONS["workers_per_org"], "Workers in dynamic high-risk environments"),
        ("year2_penetration", "Growth", "Illustrative Year 2 beachhead penetration", ASSUMPTIONS["year2_penetration"], "Base-case SOM assumption"),
        ("year3_penetration", "Growth", "Illustrative Year 3 beachhead penetration", ASSUMPTIONS["year3_penetration"], "Base-case SOM assumption"),
        ("year2_new_customers", "Growth", "Year 2 new customers", ASSUMPTIONS["year2_new_customers"], "Base-case commercialization ramp"),
        ("year3_new_customers", "Growth", "Year 3 new customers", ASSUMPTIONS["year3_new_customers"], "Base-case commercialization ramp"),
        ("year2_opex_growth", "Growth", "Year 2 OPEX growth", ASSUMPTIONS["year2_opex_growth"], "Hiring and go-to-market expansion"),
        ("year3_opex_growth", "Growth", "Year 3 OPEX growth", ASSUMPTIONS["year3_opex_growth"], "Scaling operations"),
    ]
    assumption_cells = {}
    for key, category, variable, value, note in assumption_rows:
        ws.append([category, variable, value, note])
        assumption_cells[key] = f"C{ws.max_row}"
    style_header(ws)
    autosize(ws)
    ws.freeze_panes = "A2"

    plan = wb.create_sheet("Deployment Plan")
    months = [f"M{i}" for i in range(1, 13)]
    plan.append(["Metric"] + months + ["Year 1 Total"])
    deployment_rows = {
        "New customers": DEPLOYMENT["new_customers"],
        "New hardware units": DEPLOYMENT["new_units"],
        "Active users": DEPLOYMENT["active_users"],
    }
    for label, values in deployment_rows.items():
        plan.append([label] + values + [f"=SUM(B{plan.max_row + 1}:M{plan.max_row + 1})"])
    style_header(plan)
    autosize(plan)
    plan.freeze_panes = "B2"

    pnl = wb.create_sheet("Monthly P&L")
    pnl.append(["Line Item"] + months + ["Year 1 Total"])
    line_items = [
        "Hardware revenue",
        "Software revenue",
        "Onboarding revenue",
        "Total revenue",
        "Hardware COGS",
        "Cloud cost",
        "Support cost",
        "Total COGS",
        "Gross profit",
        "Engineering payroll",
        "Sales and business development",
        "Software/tools/admin",
        "Marketing/events",
        "Total OPEX",
        "EBITDA before CAPEX",
        "CAPEX",
        "Net cash impact",
    ]
    for item in line_items:
        pnl.append([item] + [""] * 13)

    row = {pnl.cell(r, 1).value: r for r in range(2, pnl.max_row + 1)}

    for month_idx in range(2, 14):
        month_col = get_column_letter(month_idx)
        dep_col = get_column_letter(month_idx)
        pnl[f"{month_col}{row['Hardware revenue']}"] = f"='Deployment Plan'!{dep_col}3*'Assumptions'!{assumption_cells['hardware_price']}"
        pnl[f"{month_col}{row['Software revenue']}"] = f"='Deployment Plan'!{dep_col}4*'Assumptions'!{assumption_cells['software_fee']}"
        pnl[f"{month_col}{row['Onboarding revenue']}"] = f"='Deployment Plan'!{dep_col}2*'Assumptions'!{assumption_cells['onboarding_fee']}"
        pnl[f"{month_col}{row['Total revenue']}"] = f"=SUM({month_col}{row['Hardware revenue']}:{month_col}{row['Onboarding revenue']})"
        pnl[f"{month_col}{row['Hardware COGS']}"] = f"='Deployment Plan'!{dep_col}3*'Assumptions'!{assumption_cells['hardware_cogs']}"
        pnl[f"{month_col}{row['Cloud cost']}"] = f"='Deployment Plan'!{dep_col}4*'Assumptions'!{assumption_cells['cloud_cost']}"
        pnl[f"{month_col}{row['Support cost']}"] = f"='Deployment Plan'!{dep_col}4*'Assumptions'!{assumption_cells['support_cost']}"
        pnl[f"{month_col}{row['Total COGS']}"] = f"=SUM({month_col}{row['Hardware COGS']}:{month_col}{row['Support cost']})"
        pnl[f"{month_col}{row['Gross profit']}"] = f"={month_col}{row['Total revenue']}-{month_col}{row['Total COGS']}"
        pnl[f"{month_col}{row['Engineering payroll']}"] = f"='Assumptions'!{assumption_cells['eng_monthly']}"
        pnl[f"{month_col}{row['Sales and business development']}"] = f"='Assumptions'!{assumption_cells['sales_monthly']}"
        pnl[f"{month_col}{row['Software/tools/admin']}"] = f"='Assumptions'!{assumption_cells['tools_monthly']}"
        pnl[f"{month_col}{row['Marketing/events']}"] = f"='Assumptions'!{assumption_cells['marketing_monthly']}"
        pnl[f"{month_col}{row['Total OPEX']}"] = f"=SUM({month_col}{row['Engineering payroll']}:{month_col}{row['Marketing/events']})"
        pnl[f"{month_col}{row['EBITDA before CAPEX']}"] = f"={month_col}{row['Gross profit']}-{month_col}{row['Total OPEX']}"
        pnl[f"{month_col}{row['CAPEX']}"] = ASSUMPTIONS["prototype_capex"] / 4 if month_idx <= 5 else 0
        pnl[f"{month_col}{row['Net cash impact']}"] = f"={month_col}{row['EBITDA before CAPEX']}-{month_col}{row['CAPEX']}"

    for r in range(2, pnl.max_row + 1):
        pnl[f"N{r}"] = f"=SUM(B{r}:M{r})"

    style_header(pnl)
    autosize(pnl)
    pnl.freeze_panes = "B2"

    unit = wb.create_sheet("Unit Economics")
    unit.append(["Metric", "Value", "Formula / Comment"])
    unit_rows = [
        ("Hardware gross margin per unit", f"='Assumptions'!{assumption_cells['hardware_price']}-'Assumptions'!{assumption_cells['hardware_cogs']}", "Hardware price minus hardware COGS"),
        ("Monthly contribution per active user", f"='Assumptions'!{assumption_cells['software_fee']}-'Assumptions'!{assumption_cells['cloud_cost']}-'Assumptions'!{assumption_cells['support_cost']}", "Software fee minus support and cloud"),
        ("First-year customers", "='Deployment Plan'!N2", "New customers in year 1"),
        ("First-year hardware units", "='Deployment Plan'!N3", "Total deployed pilot units"),
        ("Year-end active users", "=MROUND('Deployment Plan'!M4,1)", "Active users in month 12"),
        ("Average onboarding revenue per customer", f"='Assumptions'!{assumption_cells['onboarding_fee']}", "One-time fee"),
        ("Year 1 gross margin", "='Monthly P&L'!N10/'Monthly P&L'!N5", "Gross profit divided by total revenue"),
        ("Year-end ARR", f"='Deployment Plan'!M4*'Assumptions'!{assumption_cells['software_fee']}*12", "Annual recurring software revenue exiting year 1"),
    ]
    for row_data in unit_rows:
        unit.append(list(row_data))
    style_header(unit)
    autosize(unit)

    market = wb.create_sheet("Market Sizing")
    market.append(["Metric", "Value", "Notes"])
    market_rows = [
        ("Beachhead geography", "Western Canada industrial construction and turnaround sites", "Initial focus"),
        ("Target organizations in beachhead", f"='Assumptions'!{assumption_cells['target_orgs']}", "Planning assumption to be validated"),
        ("Average qualifying workers per target org", f"='Assumptions'!{assumption_cells['workers_per_org']}", "Workers in high-risk dynamic environments"),
        ("Total beachhead end users", "=B3*B4", "Potential wearable users"),
        ("Annual software revenue at 100% beachhead penetration", f"=B5*'Assumptions'!{assumption_cells['software_fee']}*12", "ARR only, excludes hardware"),
        ("Illustrative Year 2 SOM penetration", f"='Assumptions'!{assumption_cells['year2_penetration']}", "Base-case SOM assumption"),
        ("Illustrative Year 2 active users", "=B5*B7", "Users at Year 2 SOM"),
        ("Illustrative Year 2 ARR", f"=B8*'Assumptions'!{assumption_cells['software_fee']}*12", "Software ARR at Year 2 SOM"),
        ("Illustrative Year 3 SOM penetration", f"='Assumptions'!{assumption_cells['year3_penetration']}", "Attainable share after validation"),
        ("Illustrative Year 3 active users", "=B5*B10", "Users at Year 3 SOM"),
        ("Illustrative Year 3 ARR", f"=B11*'Assumptions'!{assumption_cells['software_fee']}*12", "Software ARR at Year 3 SOM"),
        ("Global smart PPE market size, 2025 (USD M)", f"='Assumptions'!{assumption_cells['smart_ppe_market_2025_musd']}", "Top-down market context"),
        ("Global smart PPE market size, 2033 (USD M)", f"='Assumptions'!{assumption_cells['smart_ppe_market_2033_musd']}", "Top-down market context"),
        ("Global smart PPE market CAGR", f"=(B14/B13)^(1/8)-1", "2025 to 2033 CAGR"),
        ("Global smart helmet market size, 2024 (USD M)", f"='Assumptions'!{assumption_cells['smart_helmet_market_2024_musd']}", "Top-down market context"),
        ("Global smart helmet market size, 2030 (USD M)", f"='Assumptions'!{assumption_cells['smart_helmet_market_2030_musd']}", "Top-down market context"),
        ("Global smart helmet market CAGR", f"=(B17/B16)^(1/6)-1", "2024 to 2030 CAGR"),
    ]
    for row_data in market_rows:
        market.append(list(row_data))
    style_header(market)
    autosize(market)

    annual = wb.create_sheet("Annual Outlook")
    annual.append(["Metric", "Year 1", "Year 2", "Year 3"])
    annual_rows = [
        ("New customers", "='Deployment Plan'!N2", f"='Assumptions'!{assumption_cells['year2_new_customers']}", f"='Assumptions'!{assumption_cells['year3_new_customers']}"),
        ("New hardware units", "='Deployment Plan'!N3", "=ROUND('Market Sizing'!B8-'Deployment Plan'!M4,0)", "=ROUND('Market Sizing'!B11-'Market Sizing'!B8,0)"),
        ("Year-end active users", "='Deployment Plan'!M4", "=ROUND('Market Sizing'!B8,0)", "=ROUND('Market Sizing'!B11,0)"),
        ("Average active users", "='Deployment Plan'!N4/12", "=(B4+C4)/2", "=(C4+D4)/2"),
        ("Hardware revenue", "='Monthly P&L'!N2", f"=C2*'Assumptions'!{assumption_cells['hardware_price']}", f"=D2*'Assumptions'!{assumption_cells['hardware_price']}"),
        ("Software revenue", "='Monthly P&L'!N3", f"=C4*12*'Assumptions'!{assumption_cells['software_fee']}", f"=D4*12*'Assumptions'!{assumption_cells['software_fee']}"),
        ("Onboarding revenue", "='Monthly P&L'!N4", f"=C2*'Assumptions'!{assumption_cells['onboarding_fee']}", f"=D2*'Assumptions'!{assumption_cells['onboarding_fee']}"),
        ("Total revenue", "='Monthly P&L'!N5", "=SUM(C6:C8)", "=SUM(D6:D8)"),
        ("Total COGS", "='Monthly P&L'!N9", f"=(C2*'Assumptions'!{assumption_cells['hardware_cogs']})+(C4*12*('Assumptions'!{assumption_cells['cloud_cost']}+'Assumptions'!{assumption_cells['support_cost']}))", f"=(D2*'Assumptions'!{assumption_cells['hardware_cogs']})+(D4*12*('Assumptions'!{assumption_cells['cloud_cost']}+'Assumptions'!{assumption_cells['support_cost']}))"),
        ("Gross profit", "='Monthly P&L'!N10", "=C9-C10", "=D9-D10"),
        ("Gross margin", "='Monthly P&L'!N10/'Monthly P&L'!N5", "=C11/C9", "=D11/D9"),
        ("OPEX", "='Monthly P&L'!N15", f"=B13*(1+'Assumptions'!{assumption_cells['year2_opex_growth']})", f"=C13*(1+'Assumptions'!{assumption_cells['year3_opex_growth']})"),
        ("EBITDA before CAPEX", "='Monthly P&L'!N16", "=C11-C13", "=D11-D13"),
        ("Revenue growth", "", "=C9/B9-1", "=D9/C9-1"),
        ("Year-end ARR", f"='Deployment Plan'!M4*'Assumptions'!{assumption_cells['software_fee']}*12", f"=C3*'Assumptions'!{assumption_cells['software_fee']}*12", f"=D3*'Assumptions'!{assumption_cells['software_fee']}*12"),
    ]
    for row_data in annual_rows:
        annual.append(list(row_data))
    style_header(annual)
    autosize(annual)

    dashboard["A1"] = "SiteSentry Financial Dashboard"
    dashboard["A1"].font = Font(size=16, bold=True)
    dashboard["A3"] = "Market"
    dashboard["D3"] = "Beachhead"
    dashboard["A10"] = "Year 1"
    dashboard["D10"] = "3-Year Outlook"
    for cell in ["A3", "D3", "A10", "D10"]:
        dashboard[cell].font = Font(bold=True)
        dashboard[cell].fill = PatternFill("solid", fgColor="D9EAF7")

    dashboard_rows = [
        ("A4", "Global smart PPE market", f"='Market Sizing'!B13&\"M\""),
        ("A5", "Smart PPE CAGR", "='Market Sizing'!B15"),
        ("A6", "Global smart helmet market", f"='Market Sizing'!B16&\"M\""),
        ("A7", "Smart helmet CAGR", "='Market Sizing'!B18"),
        ("D4", "Beachhead end users", "='Market Sizing'!B5"),
        ("D5", "100% beachhead software ARR", "='Market Sizing'!B6"),
        ("D6", "Year 3 SOM active users", "='Market Sizing'!B11"),
        ("D7", "Year 3 ARR", "='Annual Outlook'!D16"),
        ("A11", "Year 1 revenue", "='Annual Outlook'!B9"),
        ("A12", "Year 1 gross margin", "='Annual Outlook'!B12"),
        ("A13", "Year-end ARR", "='Annual Outlook'!B16"),
        ("A14", "Net cash impact", "='Monthly P&L'!N18"),
        ("D11", "Year 2 revenue", "='Annual Outlook'!C9"),
        ("D12", "Year 2 revenue growth", "='Annual Outlook'!C15"),
        ("D13", "Year 3 revenue", "='Annual Outlook'!D9"),
        ("D14", "Year 3 gross margin", "='Annual Outlook'!D12"),
    ]
    for cell, label, formula in dashboard_rows:
        dashboard[cell] = label
        dashboard[f"{get_column_letter(dashboard[cell].column + 1)}{dashboard[cell].row}"] = formula

    revenue_chart = BarChart()
    revenue_chart.title = "Revenue Outlook"
    revenue_chart.y_axis.title = "CAD"
    data = Reference(annual, min_col=2, max_col=4, min_row=9, max_row=9)
    cats = Reference(annual, min_col=2, max_col=4, min_row=1, max_row=1)
    revenue_chart.add_data(data, titles_from_data=False)
    revenue_chart.set_categories(cats)
    revenue_chart.height = 6
    revenue_chart.width = 10
    dashboard.add_chart(revenue_chart, "A17")

    users_chart = LineChart()
    users_chart.title = "Active User Growth"
    users_chart.y_axis.title = "Users"
    users_data = Reference(annual, min_col=2, max_col=4, min_row=4, max_row=4)
    users_chart.add_data(users_data, titles_from_data=False)
    users_chart.set_categories(cats)
    users_chart.height = 6
    users_chart.width = 10
    dashboard.add_chart(users_chart, "K17")

    autosize(dashboard)

    milestones = wb.create_sheet("Milestones")
    milestones.append(["Year", "Milestones"])
    milestone_rows = [
        ("Year 1", "Build MVP, secure 4 paid pilots, validate one or two hazard classes, collect field data"),
        ("Year 2", "Convert pilots to annual contracts, improve product reliability, establish reference customers"),
        ("Year 3", "Expand into mining or manufacturing, add channel partners, grow recurring software base"),
        ("Years 4-5", "Broaden supported hazard library, integrate with enterprise safety workflows, win multi-site deployments"),
    ]
    for row_data in milestone_rows:
        milestones.append(list(row_data))
    style_header(milestones)
    autosize(milestones)

    for sheet in wb.worksheets:
        for row_cells in sheet.iter_rows():
            for cell in row_cells:
                if isinstance(cell.value, str) and cell.value.startswith("="):
                    label = str(sheet.cell(cell.row, 1).value).lower() if cell.column > 1 else ""
                    if any(term in label for term in ["revenue", "profit", "cogs", "arr", "cash", "opex"]):
                        cell.number_format = '"C$"#,##0'
                    elif "margin" in label or "growth" in label or "penetration" in label or "cagr" in label:
                        cell.number_format = "0.0%"

    for cell_ref in ["E4", "E5", "E7", "E12", "E13", "E14", "E15"]:
        dashboard[cell_ref].number_format = '"C$"#,##0'
    for cell_ref in ["B5", "B7", "E12", "E14"]:
        dashboard[cell_ref].number_format = "0.0%"

    wb.save(XLSX_PATH)
    return summary


def add_bullets(slide, title, bullets):
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0


def build_pitch_deck(summary):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "SiteSentry"
    title_slide.placeholders[1].text = "AI-enabled wearable hazard detection for industrial workers"

    slides = [
        ("Problem and Why Now", [
            "Workers in dynamic industrial environments often miss hazards outside their immediate field of view.",
            "Passive PPE reduces injury severity but does not actively improve real-time situational awareness.",
            "Near misses and incidents create human harm, delay, liability, and operating cost, and edge AI now makes active protection feasible."
        ]),
        ("Solution and Architecture", [
            "Smart hardhat platform with worker-view sensing and real-time hazard alerts.",
            "Camera and alert interface live on the hardhat, while battery and embedded compute sit in a belt-mounted module.",
            "A rugged cable reduces headborne weight and heat while preserving local low-latency inference and offline resilience."
        ]),
        ("Market, User, and Customer", [
            f"Initial focus: Western Canadian industrial construction and turnaround sites with about {summary['beachhead_users']:,} target workers across roughly {ASSUMPTIONS['target_orgs']} organizations.",
            f"Top-down market context: smart PPE grows from US${ASSUMPTIONS['smart_ppe_market_2025_musd']/1000:.1f}B in 2025 to US${ASSUMPTIONS['smart_ppe_market_2033_musd']/1000:.1f}B by 2033, about {percent(summary['ppe_cagr'])} CAGR.",
            f"End user: frontline industrial worker. Customer: employer, EHS leader, or major contractor."
        ]),
        ("Value Proposition and Competitive Position", [
            "SiteSentry adds an active safety layer that detects hazards and alerts workers before incidents happen.",
            "Employer value: avoided incidents, stronger onboarding, better safety visibility, and lower disruption.",
            "Differentiation: worker-view hazard detection that complements passive PPE, spotters, and site-level analytics."
        ]),
        ("Business Model", [
            "B2B model: hardware sale plus recurring software subscription and onboarding fee.",
            "Direct sales through pilots with large contractors and site operators.",
            "Longer term: analytics add-ons, channel partners, and multi-site enterprise contracts."
        ]),
        ("MVP and Validation", [
            "Narrow MVP focused on one form factor and one or two hazard classes.",
            "Key questions: worker acceptance, alert accuracy, buyer willingness to pay, measurable pilot value.",
            "Use paid pilots to collect evidence before broadening the feature set."
        ]),
        ("Financial Snapshot", [
            f"Year 1 base case: {summary['y1']['new_customers']} paid pilots, {summary['y1']['units']} units, {money(summary['y1']['revenue'])} revenue, and {percent(summary['y1']['gross_margin'])} gross margin.",
            f"Illustrative scale path: {money(summary['y2']['revenue'])} revenue in Year 2 and {money(summary['y3']['revenue'])} in Year 3, with gross margin expanding to {percent(summary['y3']['gross_margin'])}.",
            f"Recurring base builds from {money(summary['y1']['arr'])} year-end ARR in Year 1 to {money(summary['y3']['arr'])} by Year 3 at a {percent(ASSUMPTIONS['year3_penetration'])} beachhead SOM case."
        ]),
        ("Key Risks", [
            "Technical: false positives or poor performance in real worksites.",
            "Commercial: slow enterprise sales cycles and unclear budget ownership.",
            "Adoption and policy: comfort, privacy, recording concerns, and worker trust."
        ]),
        ("Next Steps", [
            "Run customer interviews across workers, supervisors, and buyers.",
            "Define the first pilot use case and build an MVP around it.",
            "Secure a design partner and validate whether the product delivers trusted, measurable value."
        ]),
    ]

    for title, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        add_bullets(slide, title, bullets)

    prs.save(PPTX_PATH)


if __name__ == "__main__":
    summary = build_financial_model()
    build_pitch_deck(summary)
